import os
import uuid
import logging
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("export.ppt")

# Dark Theme Colors
BG_COLOR = RGBColor(15, 23, 42)      # Slate 900
TEXT_MAIN = RGBColor(248, 250, 252)  # Slate 50
TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
ACCENT = RGBColor(59, 130, 246)      # Blue 500

class PPTExporter:
    @staticmethod
    async def generate_pitch_deck(report_data: Dict[str, Any]) -> str:
        """
        Generates an 18-slide investor-grade Pitch Deck PPTX.
        Returns the file path to the generated PPTX.
        """
        logger.info(f"Generating Pitch Deck for idea: {report_data.get('startup_idea', 'Unknown')}")
        
        output_dir = os.path.join(os.getcwd(), "reports")
        os.makedirs(output_dir, exist_ok=True)
        filename = f"venturelens_deck_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(output_dir, filename)
        
        prs = Presentation()
        # Set to 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        analysis = report_data.get('analysis_payload', {})
        idea = str(report_data.get('startup_idea', 'New Venture'))
        
        # Agents data
        market = analysis.get('market_agent', {})
        customer = analysis.get('customer_agent', {})
        comp = analysis.get('competitor_agent', {})
        gtm = analysis.get('gtm_agent', {})
        score = analysis.get('startup_score_agent', {})
        
        # --- Helper Functions ---
        def add_slide_with_bg():
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR
            return slide
            
        def add_title(slide, text, top=0.5):
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = str(text)
            p.font.bold = True
            p.font.size = Pt(36)
            p.font.color.rgb = TEXT_MAIN
            
        def add_subtitle(slide, text, top=1.3):
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(0.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = str(text)
            p.font.size = Pt(18)
            p.font.color.rgb = ACCENT
            
        def add_bullets(slide, items, left=0.8, top=2.0, width=11.5, height=4.5, font_size=24):
            if not items: return
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.add_paragraph()
                p.text = str(item)
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT_MUTED
                p.level = 0
                p.space_after = Pt(14)
                
        def add_stat_box(slide, left, top, value, label):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.5), Inches(2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
            shape.line.color.rgb = ACCENT
            
            tf = shape.text_frame
            p1 = tf.add_paragraph()
            p1.text = str(value)
            p1.font.bold = True
            p1.font.size = Pt(32)
            p1.font.color.rgb = TEXT_MAIN
            p1.alignment = PP_ALIGN.CENTER
            
            p2 = tf.add_paragraph()
            p2.text = str(label)
            p2.font.size = Pt(14)
            p2.font.color.rgb = TEXT_MUTED
            p2.alignment = PP_ALIGN.CENTER
            
        def safe_extract(d, key, default):
            if not d: return default
            val = d.get(key)
            if not val or val == "N/A" or val == "None" or val == "": return default
            return val

        # 1: Cover
        s1 = add_slide_with_bg()
        add_title(s1, idea[:100] + "..." if len(idea)>100 else idea, top=2.5)
        add_subtitle(s1, "Confidential Investor Pitch Deck", top=3.5)
        
        # 2: Problem Statement
        s2 = add_slide_with_bg()
        add_title(s2, "The Problem")
        add_subtitle(s2, "Why does this need to exist?")
        pain_points = customer.get("pain_points", ["Inefficient legacy processes causing time and money loss", "Lack of specialized tooling for modern workflows"])
        add_bullets(s2, pain_points[:4] if isinstance(pain_points, list) else [pain_points])
        
        # 3: Market Opportunity
        s3 = add_slide_with_bg()
        add_title(s3, "Market Opportunity")
        add_subtitle(s3, "A rapidly expanding market ripe for disruption")
        add_stat_box(s3, 1.0, 2.5, safe_extract(market, "market_size", "Multi-Billion"), "Total Addressable Market")
        add_stat_box(s3, 5.0, 2.5, safe_extract(market, "growth_rate", "High CAGR"), "Expected Growth Rate")
        add_stat_box(s3, 9.0, 2.5, "High", "Market Urgency")
        
        # 4: Solution
        s4 = add_slide_with_bg()
        add_title(s4, "Our Solution")
        add_subtitle(s4, "A paradigm shift in how the problem is solved")
        add_bullets(s4, ["Purpose-built to address core pain points directly", "Automates manual workflows to 10x productivity", "Seamless integration into existing user habits"])
        
        # 5: Product Demo / MVP
        s5 = add_slide_with_bg()
        add_title(s5, "The Product (MVP)")
        add_subtitle(s5, "Core value delivered on day one")
        features = gtm.get("mvp_features", ["Core workflow automation", "Intuitive user dashboard", "Data export and reporting"])
        add_bullets(s5, [f.get("feature", "Core Feature") if isinstance(f, dict) else f for f in features][:4])
        
        # 6: Why Now
        s6 = add_slide_with_bg()
        add_title(s6, "Why Now?")
        add_subtitle(s6, "Timing is everything in venture")
        trends = market.get("market_trends", ["Rapid technology adoption in legacy sectors", "Shifting consumer expectations"])
        add_bullets(s6, trends[:3] if isinstance(trends, list) else [trends])

        # 7: TAM / SAM / SOM
        s7 = add_slide_with_bg()
        add_title(s7, "TAM / SAM / SOM")
        add_subtitle(s7, "Bottom-up market sizing")
        add_stat_box(s7, 1.0, 2.5, safe_extract(market, "market_size", "Large"), "TAM (Total Market)")
        add_stat_box(s7, 5.0, 2.5, "30% of TAM", "SAM (Serviceable Market)")
        add_stat_box(s7, 9.0, 2.5, "1-5% of SAM", "SOM (Obtainable Market)")
        
        # 8: Customer Personas
        s8 = add_slide_with_bg()
        add_title(s8, "Target Customer")
        add_subtitle(s8, "Who feels the pain most acutely?")
        segments = customer.get("target_customer_segments", ["Enterprise Decision Makers"])
        add_bullets(s8, segments[:3] if isinstance(segments, list) else [segments])
            
        # 9: Competitive Landscape
        s9 = add_slide_with_bg()
        add_title(s9, "Competitive Landscape")
        add_subtitle(s9, "How we map against incumbents")
        add_bullets(s9, ["vs. Status Quo: Faster and more automated", "vs. Point Solutions: Unified platform experience", "vs. Legacy Enterprise: Modern, cost-effective, and agile"])
            
        # 10: Unique Advantage / Moat
        s10 = add_slide_with_bg()
        add_title(s10, "Our Moat")
        add_subtitle(s10, "Sustainable competitive advantage")
        add_bullets(s10, [safe_extract(comp, "gap_analysis", "Leveraging automation to eliminate friction."), "First-mover advantage in a niche vertical", "Proprietary data aggregation loop"])
        
        # 11: Business Model
        s11 = add_slide_with_bg()
        add_title(s11, "Business Model")
        add_subtitle(s11, "How we make money")
        add_bullets(s11, [safe_extract(gtm, "business_model", "Tiered B2B SaaS subscription model."), "Predictable recurring revenue (ARR)", "High gross margins typical of SaaS"])
        
        # 12: Go-To-Market Strategy
        s12 = add_slide_with_bg()
        add_title(s12, "Go-To-Market")
        add_subtitle(s12, "Acquiring our first 1,000 customers")
        channels = gtm.get("acquisition_channels", ["Direct Sales", "Content Marketing", "Strategic Partnerships"])
        add_bullets(s12, [c.get("channel", str(c)) if isinstance(c, dict) else str(c) for c in channels][:3])
        
        # 13: Traction & Validation
        s13 = add_slide_with_bg()
        add_title(s13, "Validation & Traction")
        add_subtitle(s13, "AI-driven viability scoring")
        add_stat_box(s13, 1.0, 2.5, f"{score.get('overall_score', 85)}/100", "AI Validation Score")
        add_stat_box(s13, 5.0, 2.5, "Validated", "Problem-Solution Fit")
        add_stat_box(s13, 9.0, 2.5, score.get("confidence_level", "High"), "AI Confidence")
        
        # 14: Product Roadmap
        s14 = add_slide_with_bg()
        add_title(s14, "Product Roadmap")
        add_subtitle(s14, "The next 12-18 months")
        add_stat_box(s14, 1.0, 2.5, "Q1", "MVP Launch & Core Validation")
        add_stat_box(s14, 5.0, 2.5, "Q2-Q3", "Expansion of Feature Set")
        add_stat_box(s14, 9.0, 2.5, "Q4", "Scale & Enterprise Tier")
        
        # 15: Financial Projections
        s15 = add_slide_with_bg()
        add_title(s15, "Financial Horizon")
        add_subtitle(s15, "Path to profitability")
        add_bullets(s15, ["Year 1: Focus on product-market fit and early MRR", "Year 2: Scalable acquisition and channel expansion", "Year 3: Hitting inflection point for Series A/B metrics"])
        
        # 16: Team
        s16 = add_slide_with_bg()
        add_title(s16, "The Team")
        add_subtitle(s16, "Why we are the right people to build this")
        add_bullets(s16, ["Deep domain expertise in the target industry", "Technical capabilities to execute the MVP rapidly", "Obsession with solving the customer's problem"])
        
        # 17: Funding Ask
        s17 = add_slide_with_bg()
        add_title(s17, "The Ask")
        add_subtitle(s17, "Fueling the next stage of growth")
        add_bullets(s17, ["Raising Pre-Seed / Seed round", "Capital allocation: 50% Engineering, 30% GTM/Sales, 20% Operations", "Milestone: Achieve $1M ARR within 18 months"])
        
        # 18: Closing Vision
        s18 = add_slide_with_bg()
        add_title(s18, "Join Us.", top=3.0)
        add_subtitle(s18, score.get("verdict", "Building the future of this industry."), top=4.0)

        prs.save(filepath)
        logger.info(f"Successfully generated investor deck: {filename}")
        return filepath
